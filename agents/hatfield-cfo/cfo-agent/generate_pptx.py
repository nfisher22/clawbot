#!/usr/bin/env python3
"""Generate CFO Presentation PowerPoint for Roxster, Ltd."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x1F, 0x38, 0x64)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MED_GRAY = RGBColor(0xD9, 0xD9, 0xD9)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
AMBER = RGBColor(0xC0, 0x6E, 0x00)

FONT_NAME = "Calibri"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]


def add_shape(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    return shape


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT, font_name=FONT_NAME):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return p


def add_bullet_para(tf, text, size=14, bold=False, color=BLACK, level=0, space_before=Pt(4), space_after=Pt(2)):
    p = tf.add_paragraph()
    p.level = level
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT_NAME
    return p


def add_title_bar(slide, title_text, top=Inches(0), height=Inches(0.9)):
    bar = add_shape(slide, Inches(0), top, SLIDE_W, height, NAVY)
    bar.text_frame.word_wrap = True
    bar.text_frame.margin_left = Inches(0.6)
    bar.text_frame.margin_top = Inches(0.15)
    set_text(bar.text_frame, title_text, size=28, bold=True, color=WHITE)
    return bar


def add_footer(slide, text="F & W Properties Inc. | Peak 10 Group | Confidential"):
    tb = add_textbox(slide, Inches(0), Inches(7.05), SLIDE_W, Inches(0.4))
    tf = tb.text_frame
    set_text(tf, text, size=9, color=RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.CENTER)


def add_table(slide, rows, cols, left, top, width, height):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table_shape.table


def style_header_row(table, col_count, texts, font_size=11):
    for i in range(col_count):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT
        run = p.add_run()
        run.text = texts[i]
        run.font.size = Pt(font_size)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = FONT_NAME
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.08)


def style_data_cell(table, row, col, text, bold=False, align=None, color=BLACK, font_size=11):
    cell = table.cell(row, col)
    if row % 2 == 0:
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_GRAY
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = cell.text_frame.paragraphs[0]
    if align is None:
        p.alignment = PP_ALIGN.RIGHT if col > 0 else PP_ALIGN.LEFT
    else:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT_NAME
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)


# ==================== SLIDE 1 — TITLE ====================
slide = prs.slides.add_slide(blank_layout)
# Full navy background
add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)

# Company name
tb = add_textbox(slide, Inches(1), Inches(1.2), Inches(11.3), Inches(0.8))
set_text(tb.text_frame, "F & W Properties Inc.", size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Divider line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(2.1), Inches(4.3), Inches(0.03))
line.fill.solid()
line.fill.fore_color.rgb = WHITE
line.line.fill.background()

# Subtitle
tb = add_textbox(slide, Inches(1), Inches(2.4), Inches(11.3), Inches(0.7))
set_text(tb.text_frame, "2026 Budget & Valuation Analysis", size=30, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

# Property
tb = add_textbox(slide, Inches(1), Inches(3.3), Inches(11.3), Inches(0.6))
set_text(tb.text_frame, "Roxster, Ltd. — Roxbury Arms", size=22, bold=False, color=RGBColor(0xBF, 0xBF, 0xBF), align=PP_ALIGN.CENTER)

# Address
tb = add_textbox(slide, Inches(1), Inches(3.9), Inches(11.3), Inches(0.5))
set_text(tb.text_frame, "1571-1619 Roxbury Road, Grandview Heights, OH 43212", size=14, color=RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.CENTER)

# Prepared by
tb = add_textbox(slide, Inches(1), Inches(5.0), Inches(11.3), Inches(0.5))
set_text(tb.text_frame, "Prepared by Mr. Soul CFO | March 17, 2026", size=16, color=RGBColor(0xBF, 0xBF, 0xBF), align=PP_ALIGN.CENTER)

# Confidential
tb = add_textbox(slide, Inches(1), Inches(5.6), Inches(11.3), Inches(0.4))
set_text(tb.text_frame, "Peak 10 Group | Confidential", size=12, color=RGBColor(0x80, 0x80, 0x80), align=PP_ALIGN.CENTER)


# ==================== SLIDE 2 — PURPOSE & SCOPE ====================
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Purpose & Scope")

tb = add_textbox(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(5.5))
tf = tb.text_frame
tf.word_wrap = True

bullets = [
    "2026 operating budget for Roxster, Ltd. (Roxbury Arms)",
    "Analysis covers NOI, debt service, and cash flow projections",
    "Income-approach valuation at 7.5% cap rate",
    "Sensitivity analysis across cap rates 6.0%–8.5%",
    "Key risks: payroll anomaly in 2025, elevated R&M, heavy CapEx front-loading",
    "Decision context: budget approval, CapEx authorization, valuation next steps",
]

set_text(tf, bullets[0], size=18, color=DARK_TEXT)
tf.paragraphs[0].space_before = Pt(12)
tf.paragraphs[0].space_after = Pt(8)

for b in bullets[1:]:
    add_bullet_para(tf, b, size=18, color=DARK_TEXT, space_before=Pt(12), space_after=Pt(8))

add_footer(slide)


# ==================== SLIDE 3 — PROPERTY SUMMARY ====================
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Property Summary — Roxster, Ltd.")

data = [
    ("Total Operating Income", "$1,335,287"),
    ("Total Operating Expense (Adj.)", "$706,633"),
    ("Net Operating Income (Adj.)", "$628,654"),
    ("Mortgage Interest", "$219,199"),
    ("Principal Reduction", "$23,468"),
    ("Cash Flow Before CapEx", "$385,987"),
    ("Capital Expenditures", "$341,000"),
    ("Cash Flow After CapEx", "$44,987"),
    ("Estimated Value (7.5% Cap)", "$8,382,055"),
    ("Estimated Equity Residual", "$2,930,593"),
]

table = add_table(slide, len(data) + 1, 2, Inches(2), Inches(1.2), Inches(9.3), Inches(5.5))
style_header_row(table, 2, ["Metric", "Amount"], font_size=13)

for i, (metric, amount) in enumerate(data):
    r = i + 1
    bold = metric in ("Net Operating Income (Adj.)", "Cash Flow After CapEx", "Estimated Value (7.5% Cap)", "Estimated Equity Residual")
    style_data_cell(table, r, 0, metric, bold=bold, font_size=12)
    style_data_cell(table, r, 1, amount, bold=bold, font_size=12)

table.columns[0].width = Inches(5.5)
table.columns[1].width = Inches(3.8)

add_footer(slide)


# ==================== SLIDE 4 — INCOME & EXPENSE DETAIL ====================
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Roxster, Ltd. — Income & Expense Detail")

# INCOME table
inc_label = add_textbox(slide, Inches(0.5), Inches(1.05), Inches(3), Inches(0.35))
set_text(inc_label.text_frame, "INCOME", size=13, bold=True, color=NAVY)

inc_data = [
    ("Rents", "$1,274,267"),
    ("Fee Income", "$10,356"),
    ("Other Income", "$9,264"),
    ("Utility Reimbursement", "$41,400"),
    ("Total Operating Income", "$1,335,287"),
]

t1 = add_table(slide, len(inc_data) + 1, 2, Inches(0.5), Inches(1.35), Inches(5.5), Inches(2.2))
style_header_row(t1, 2, ["Category", "Annual Budget"], font_size=10)
for i, (cat, amt) in enumerate(inc_data):
    bold = cat.startswith("Total")
    style_data_cell(t1, i+1, 0, cat, bold=bold, font_size=10)
    style_data_cell(t1, i+1, 1, amt, bold=bold, font_size=10)
t1.columns[0].width = Inches(3.3)
t1.columns[1].width = Inches(2.2)

# EXPENSE table
exp_label = add_textbox(slide, Inches(0.5), Inches(3.65), Inches(3), Inches(0.35))
set_text(exp_label.text_frame, "EXPENSES", size=13, bold=True, color=NAVY)

exp_data = [
    ("Administrative", "$37,473"),
    ("Payroll", "$265,055"),
    ("Utilities", "$74,088"),
    ("Repairs & Maintenance", "$168,516"),
    ("Taxes & Insurance", "$101,388"),
    ("Professional Services", "$76,157"),
    ("Total Operating Expense", "$706,633"),
]

t2 = add_table(slide, len(exp_data) + 1, 2, Inches(0.5), Inches(3.95), Inches(5.5), Inches(2.7))
style_header_row(t2, 2, ["Category", "Annual Budget (Adj.)"], font_size=10)
for i, (cat, amt) in enumerate(exp_data):
    bold = cat.startswith("Total")
    style_data_cell(t2, i+1, 0, cat, bold=bold, font_size=10)
    style_data_cell(t2, i+1, 1, amt, bold=bold, font_size=10)
t2.columns[0].width = Inches(3.3)
t2.columns[1].width = Inches(2.2)

# KEY METRICS on the right
km_label = add_textbox(slide, Inches(6.8), Inches(1.05), Inches(3), Inches(0.35))
set_text(km_label.text_frame, "KEY METRICS", size=13, bold=True, color=NAVY)

km_data = [
    ("NOI (Adjusted)", "$628,654"),
    ("Mortgage Interest", "$219,199"),
    ("Principal Reduction", "$23,468"),
    ("Cash Flow Before CapEx", "$385,987"),
    ("Cash Flow After CapEx", "$44,987"),
]

t3 = add_table(slide, len(km_data) + 1, 2, Inches(6.8), Inches(1.35), Inches(5.8), Inches(2.2))
style_header_row(t3, 2, ["Metric", "Amount"], font_size=10)
for i, (met, amt) in enumerate(km_data):
    bold = met.startswith("Cash Flow After")
    style_data_cell(t3, i+1, 0, met, bold=bold, font_size=10)
    style_data_cell(t3, i+1, 1, amt, bold=bold, font_size=10)
t3.columns[0].width = Inches(3.5)
t3.columns[1].width = Inches(2.3)

# Warning note
warn_tb = add_textbox(slide, Inches(6.8), Inches(3.75), Inches(5.8), Inches(0.8))
warn_tf = warn_tb.text_frame
warn_tf.word_wrap = True
set_text(warn_tf, "Cash flow after CapEx is positive at $44,987 but thin — $341K CapEx concentrated in H1", size=11, bold=False, color=AMBER)

add_footer(slide)


# ==================== SLIDE 5 — VALUATION SUMMARY ====================
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Valuation Summary")

t4 = add_table(slide, 2, 5, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.9))
style_header_row(t4, 5, ["NOI", "Cap Rate", "Est. Value", "Outstanding Debt", "Equity Residual"], font_size=12)
vals = ["$628,654", "7.50%", "$8,382,055", "$5,200,000", "$2,930,593"]
for i, v in enumerate(vals):
    style_data_cell(t4, 1, i, v, bold=True, font_size=13, align=PP_ALIGN.CENTER)

t4.columns[0].width = Inches(2.2)
t4.columns[1].width = Inches(1.8)
t4.columns[2].width = Inches(2.8)
t4.columns[3].width = Inches(2.8)
t4.columns[4].width = Inches(2.7)

# Key observations
obs_label = add_textbox(slide, Inches(0.6), Inches(2.6), Inches(5), Inches(0.4))
set_text(obs_label.text_frame, "Key Observations", size=18, bold=True, color=NAVY)

obs_tb = add_textbox(slide, Inches(0.6), Inches(3.1), Inches(12), Inches(3.5))
obs_tf = obs_tb.text_frame
obs_tf.word_wrap = True

obs_bullets = [
    "Property value estimated at $8.4M at 7.5% cap rate",
    "Equity residual of $2.9M after debt and 3% closing costs",
    "50bp cap rate compression adds ~$449K to value",
    "Third-party appraisal recommended before transaction decisions",
]

set_text(obs_tf, obs_bullets[0], size=16, color=DARK_TEXT)
obs_tf.paragraphs[0].space_before = Pt(8)
obs_tf.paragraphs[0].space_after = Pt(6)
for b in obs_bullets[1:]:
    add_bullet_para(obs_tf, b, size=16, color=DARK_TEXT, space_before=Pt(8), space_after=Pt(6))

add_footer(slide)


# ==================== SLIDE 6 — OPERATING METRICS ====================
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Operating Metrics")

metrics = [
    ("Expense Ratio (Adjusted)", "52.9%"),
    ("NOI Margin", "47.1%"),
    ("Debt Service Coverage Ratio", "2.59x"),
]

t5 = add_table(slide, len(metrics) + 1, 2, Inches(2.5), Inches(1.4), Inches(8.3), Inches(1.8))
style_header_row(t5, 2, ["Metric", "Value"], font_size=14)
for i, (m, v) in enumerate(metrics):
    style_data_cell(t5, i+1, 0, m, bold=True, font_size=14)
    style_data_cell(t5, i+1, 1, v, bold=True, font_size=14, align=PP_ALIGN.CENTER)
t5.columns[0].width = Inches(5.5)
t5.columns[1].width = Inches(2.8)

# Key takeaway box
box = add_shape(slide, Inches(2.5), Inches(3.6), Inches(8.3), Inches(1.5), LIGHT_GRAY)
box.text_frame.word_wrap = True
box.text_frame.margin_left = Inches(0.3)
box.text_frame.margin_top = Inches(0.2)
box.text_frame.margin_right = Inches(0.3)
p = box.text_frame.paragraphs[0]
run = p.add_run()
run.text = "Key Takeaway"
run.font.size = Pt(16)
run.font.bold = True
run.font.color.rgb = NAVY
run.font.name = FONT_NAME

p2 = box.text_frame.add_paragraph()
p2.space_before = Pt(8)
run2 = p2.add_run()
run2.text = "Strong DSCR of 2.59x provides comfortable debt service coverage. Expense ratio slightly above the 50% target at 52.9%. NOI margin remains healthy at 47.1%."
run2.font.size = Pt(14)
run2.font.color.rgb = DARK_TEXT
run2.font.name = FONT_NAME

add_footer(slide)


# ==================== SLIDE 7 — KEY TAKEAWAYS ====================
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Key Takeaways")

takeaways = [
    ("$628,654 Adjusted NOI", "1.0% above 2025 actual NOI of $622,397"),
    ("$8.38M Estimated Value", "Income approach at 7.5% cap rate"),
    ("$2.93M Equity Residual", "After $5.2M debt and 3% closing costs"),
    ("$385,987 Cash Flow Before CapEx", "Strong debt service coverage of 2.59x"),
    ("$341,000 CapEx Budget", "Front-loaded in H1; cash reserves must support Q1–Q2 draws"),
    ("Payroll Anomaly", "March 2025 payroll of $6,091 vs $23,012 budget requires investigation"),
]

y_pos = Inches(1.15)
for headline, detail in takeaways:
    tb = add_textbox(slide, Inches(0.8), y_pos, Inches(11.7), Inches(0.35))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = headline
    r1.font.size = Pt(15)
    r1.font.bold = True
    r1.font.color.rgb = NAVY
    r1.font.name = FONT_NAME

    r2 = p.add_run()
    r2.text = "  —  " + detail
    r2.font.size = Pt(14)
    r2.font.bold = False
    r2.font.color.rgb = DARK_TEXT
    r2.font.name = FONT_NAME

    y_pos += Inches(0.85)

# Warning highlight for last item
warn_box = add_shape(slide, Inches(0.6), y_pos - Inches(0.25), Inches(12), Inches(0.03), AMBER)

add_footer(slide)


# ==================== SLIDE 8 — CONCLUSION & NEXT STEPS ====================
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Conclusion & Next Steps")

steps = [
    "Approve 2026 operating budget for Roxster, Ltd.",
    "Investigate March 2025 payroll anomaly and confirm 2026 payroll assumptions",
    "Authorize $341K CapEx plan — confirm Q1 cash reserves for $150K roofing",
    "Confirm R&M budget of $14,043/month against elevated 2025 actuals",
    "Decide on property valuation: hold, order appraisal, or adjust cap rate",
    "Review tax & insurance budget increase of 48% vs prior year",
]

tb = add_textbox(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(5.5))
tf = tb.text_frame
tf.word_wrap = True

for i, step in enumerate(steps):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(14)
    p.space_after = Pt(6)

    r_num = p.add_run()
    r_num.text = f"{i+1}.  "
    r_num.font.size = Pt(17)
    r_num.font.bold = True
    r_num.font.color.rgb = NAVY
    r_num.font.name = FONT_NAME

    r_text = p.add_run()
    r_text.text = step
    r_text.font.size = Pt(17)
    r_text.font.bold = False
    r_text.font.color.rgb = DARK_TEXT
    r_text.font.name = FONT_NAME

add_footer(slide, "F & W Properties Inc. | Peak 10 Group | Confidential")

# Save
output_path = "/Users/natefisher/Projects/hatfield/cfo-agent/output/Roxster_CFO_Presentation_Mar2026.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
